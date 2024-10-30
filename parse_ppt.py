import os
from pptx import Presentation
import base64
import urllib.parse

def get_shape_text(shape):
    text = ""
    try:
        # 그룹화된 도형 처리
        if hasattr(shape, "shape_type"):
            if shape.shape_type == 6:  # 6은 그룹 도형을 의미
                for subshape in shape.shapes:
                    text += get_shape_text(subshape)
        
        # 텍스트 처리
        if hasattr(shape, "text") and len(shape.text.strip()) > 0:
            text += shape.text + "\n"
        
        # 표 처리
        if hasattr(shape, "table"):
            if len(shape.table.rows) == 0:
                return text
                    
            # 마크다운 테이블 헤더 생성
            for cell in shape.table.rows[0].cells:
                text += f"| {cell.text.strip()} "
            text += "|\n"
            
            # 구분선 추가
            text += "|" + "---|" * len(shape.table.rows[0].cells) + "\n"
            
            # 테이블 내용 추가 (첫 번째 행을 제외하고 반복)
            for i in range(1, len(shape.table.rows)):
                for cell in shape.table.rows[i].cells:
                    text += f"| {cell.text.strip()} "
                text += "|\n"
            
            text += "\n"
        
        # 표 처리
        try:
            if shape.table:  # table 속성 존재 여부 확인
                if len(shape.table.rows) == 0:
                    return text
                    
                # 마크다운 테이블 헤더 생성
                for cell in shape.table.rows[0].cells:
                    text += f"| {cell.text.strip()} "
                text += "|\n"
                
                # 구분선 추가
                text += "|" + "---|" * len(shape.table.rows[0].cells) + "\n"
                
                # 테이블 내용 추가 (첫 번째 행을 제외하고 반복)
                for i in range(1, len(shape.table.rows)):
                    for cell in shape.table.rows[i].cells:
                        text += f"| {cell.text.strip()} "
                    text += "|\n"
                
                text += "\n"
        except (AttributeError, ValueError):
            # 테이블이 없는 경우 무시하고 진행
            pass
                    
    except Exception as e:
        print(f"처리할 수 없는 도형 유형: {type(shape)}, 오류: {str(e)}")
        
    return text

def process_pptx(file_path):
    prs = Presentation(file_path)
    slides_text = {}
    
    for slide_number, slide in enumerate(prs.slides):
        current_slide_text = []

        # 슬라이드 제목 처리
        if slide.shapes.title:
            current_slide_text.append(f"## {slide.shapes.title.text}\n")
        
        # 슬라이드 내용 처리
        for shape in slide.shapes:
            text = get_shape_text(shape)
            if text:
                current_slide_text.append(text)
        
        # 슬라이드 노트 처리
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                current_slide_text.append("\n### 슬라이드노트:\n")
                current_slide_text.append(notes_text + "\n")
        
        if current_slide_text:
            slides_text[slide_number + 1] = ''.join(current_slide_text)
    
    return slides_text

def decode_filename(filename):
    decoding_attempts = [
        lambda name: base64.b64decode(name).decode('utf-8'),
        lambda name: urllib.parse.unquote(name),
        lambda name: urllib.parse.unquote(urllib.parse.unquote(name)),
        lambda name: name.encode('iso-8859-1').decode('utf-8'),
        lambda name: name.encode('iso-8859-1').decode('utf-8')
    ]

    for attempt in decoding_attempts:
        try:
            decoded = attempt(filename)
            if any('\u3131' <= char <= '\uD79D' for char in decoded):
                return decoded
        except:
            continue

    return filename

# result 폴더가 없으면 생성
result_folder = "result"
if not os.path.exists(result_folder):
    os.makedirs(result_folder)

# data 폴더의 모든 pptx 파일 처리
data_folder = "data"
for filename in os.listdir(data_folder):
    if filename.endswith(".pptx"):
        try:
            # 파일명 디코딩 추가
            decoded_filename = decode_filename(filename)
            file_path = os.path.join(data_folder, filename)
            
            # 결과 파일명 생성 (확장자를 txt로 변경)
            output_filename = os.path.splitext(decoded_filename)[0] + "_extracted.md"
            output_path = os.path.join(result_folder, output_filename)
            
            # 파일이 이미 존재하는 경우 건너뛰기
            if os.path.exists(output_path):
                print(f"파일이 이미 존재합니다: {output_filename}")
                continue
            
            # PPT 처리
            slides_text = process_pptx(file_path)
            
            # 결과 저장
            with open(output_path, 'w', encoding='utf-8') as file:            
                for slide_num in sorted(slides_text.keys()):
                    file.write(f"# Slide {slide_num}:\n{slides_text[slide_num]}\n\n")
                    
            print(f"성공적으로 처리됨: {filename}")
            
        except Exception as e:
            print(f"파일 처리 중 오류 발생: {filename}")
            print(f"오류 내용: {str(e)}")
            continue
